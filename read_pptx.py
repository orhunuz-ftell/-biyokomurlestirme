# -*- coding: utf-8 -*-
from pptx import Presentation
import sys
import io

sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

prs = Presentation('tik4_sunum.pptx')
print(f'Toplam slayt sayisi: {len(prs.slides)}\n')

for i, slide in enumerate(prs.slides):
    print(f'\n=== SLAYT {i+1} ===')
    for shape in slide.shapes:
        if hasattr(shape, 'text') and shape.text.strip():
            print(shape.text)
            print()
